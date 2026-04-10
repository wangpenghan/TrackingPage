from pptx import Presentation
import os

ppt_path = r"d:\TraeProject\TrackingPage\.trae\documents\PPT---参考样例.pptx"
output_path = r"d:\TraeProject\TrackingPage\ppt_template_content.txt"

try:
    # 打开PPT文件
    prs = Presentation(ppt_path)
    
    output_content = []
    output_content.append("=== PPT模板结构分析 ===")
    output_content.append(f"总幻灯片数: {len(prs.slides)}")
    output_content.append("\n")

    # 遍历每一页幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        output_content.append(f"=== 第 {slide_num} 页幻灯片 ===")
        if slide.slide_layout:
            output_content.append(f"布局名称: {slide.slide_layout.name}")
        output_content.append("\n内容:")
        
        # 遍历所有形状提取文本
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                output_content.append(f"- {text}")
        
        output_content.append("\n")

    # 保存输出内容
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(output_content))
    
    print(f"PPT模板内容已提取到: {output_path}")
    print("\n提取完成！")

except Exception as e:
    print(f"提取过程发生错误: {str(e)}")
