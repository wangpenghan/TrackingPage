from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import datetime

# 模板路径和输出路径
template_path = r"d:\TraeProject\TrackingPage\.trae\documents\PPT---参考样例.pptx"
output_path = r"d:\TraeProject\TrackingPage\铁路旅客服务系统综合显示子系统技术标准申报汇报.pptx"

# 读取素材内容
with open(r"d:\TraeProject\TrackingPage\.trae\documents\素材.md", "r", encoding="utf-8") as f:
    content = f.read()

# 打开模板PPT
prs = Presentation(template_path)

# ----------------------
# 第1页：标题幻灯片
# ----------------------
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if hasattr(shape, "text"):
        if "信息化专业国铁集团技术标准" in shape.text:
            shape.text = "信息化专业国铁集团技术标准\n申请汇报"
        elif "XXXXXXX（标准名称）" in shape.text:
            shape.text = "《铁路旅客服务系统综合显示子系统技术条件》(Q/CR 663—2018)修订"
        elif "部门：XXXX事业部" in shape.text:
            shape.text = "部门：XXXX事业部         汇报人： XXXXXX"

# ----------------------
# 第2页：目录
# ----------------------
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if hasattr(shape, "text") and "目  录" in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = "目  录"
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = "\n1. 项目概况"
        p.font.size = Pt(24)
        
        p = tf.add_paragraph()
        p.text = "2. 必要性说明"
        p.font.size = Pt(24)
        
        p = tf.add_paragraph()
        p.text = "3. 可行性说明"
        p.font.size = Pt(24)
        
        p = tf.add_paragraph()
        p.text = "4. 标准章节结构及内容说明"
        p.font.size = Pt(24)
        
        p = tf.add_paragraph()
        p.text = "5. 标准制定工作进度计划"
        p.font.size = Pt(24)

# ----------------------
# 第3页：项目概况
# ----------------------
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if hasattr(shape, "text"):
        if "1 项目概况" in shape.text:
            # 替换内容
            tf = shape.text_frame
            tf.clear()
            
            p = tf.add_paragraph()
            p.text = "1 项目概况"
            p.font.size = Pt(28)
            p.font.bold = True
            
            p = tf.add_paragraph()
            p.text = "\n标准负责人："
            p.font.size = Pt(18)
            
            p = tf.add_paragraph()
            p.text = "\n标准范围："
            p.font.size = Pt(18)
            p = tf.add_paragraph()
            p.text = "本标准对铁路旅客服务系统综合显示子系统的功能要求、性能指标、接口规范、安全机制等方面作出规定。"
            p.font.size = Pt(16)
            p.level = 1
            p = tf.add_paragraph()
            p.text = "本标准适用于铁路旅客服务系统综合显示子系统的设计、开发、测试、部署和运维。"
            p.font.size = Pt(16)
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "\n预计起草单位："
            p.font.size = Pt(18)
            p = tf.add_paragraph()
            p.text = "本标准起草单位：[三家及以上单位名称]"
            p.font.size = Pt(16)
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "\n起草周期：☑ 12月 □18月 □ 24月"
            p.font.size = Pt(18)
            
            p = tf.add_paragraph()
            p.text = "\n标准化对象定义："
            p.font.size = Pt(18)
            p = tf.add_paragraph()
            p.text = "铁路旅客服务系统综合显示子系统是为旅客提供出行信息服务的显示终端及控制系统，包括引导屏、候车屏、检票屏、站台屏等各类显示设备及控制管理系统。"
            p.font.size = Pt(16)
            p.level = 1

# ----------------------
# 第4页：必要性说明 - 对应申报指南类别
# ----------------------
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if hasattr(shape, "text") and "2 必要性说明" in shape.text:
        tf = shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = "2 必要性说明"
        p.font.size = Pt(28)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "\n对应申报指南类别："
        p.font.size = Pt(18)
        p = tf.add_paragraph()
        p.text = "本项目符合《国铁集团技术标准和计量规范申报指南》中的："
        p.font.size = Pt(16)
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• 第2类：支撑国铁集团数智化发展，提升数据安全与网络安全保护能力"
        p.font.size = Pt(16)
        p.level = 2
        p = tf.add_paragraph()
        p.text = "• 第4类：提升客运服务质量，优化旅客服务体验"
        p.font.size = Pt(16)
        p.level = 2
        p = tf.add_paragraph()
        p.text = "• 第5类：推动设备升级改造，促进绿色低碳转型"
        p.font.size = Pt(16)
        p.level = 2
        p = tf.add_paragraph()
        p.text = "• 第10类：修订与新技术、新需求不匹配的技术标准"
        p.font.size = Pt(16)
        p.level = 2
        
        p = tf.add_paragraph()
        p.text = "\n使用需求/使用现状："
        p.font.size = Pt(18)
        p = tf.add_paragraph()
        p.text = "铁路旅客服务系统显示子系统已在全路各级车站广泛应用，随着Mini LED等新型显示技术的发展和网络安全要求的提升，原有标准Q/CR 663—2018已不能满足当前技术发展和业务需求，亟需修订完善。"
        p.font.size = Pt(16)
        p.level = 1

# ----------------------
# 第5页：必要性说明 - 制定标准的必要性
# ----------------------
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if hasattr(shape, "text") and "2 必要性说明" in shape.text:
        tf = shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = "2 必要性说明"
        p.font.size = Pt(28)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "\n制定标准的必要性："
        p.font.size = Pt(18)
        p = tf.add_paragraph()
        p.text = "• 现有标准缺少引导屏控制卡回读功能要求，无法实现终端设备运行状态的实时监控与远程诊断，运维效率低"
        p.font.size = Pt(16)
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• 缺乏系统通信、数据存储、信息展示等关键环节的加密与审计机制，存在数据安全和内容篡改风险"
        p.font.size = Pt(16)
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• 全彩显示屏缺少统一的信息展示版式规范，信息展示不规范、可读性差，影响旅客服务体验"
        p.font.size = Pt(16)
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• 原有标准未涵盖Mini LED等新型显示技术和智能控制单元等新设备的技术要求，兼容性和扩展性不足"
        p.font.size = Pt(16)
        p.level = 1

# ----------------------
# 第6页：可行性说明
# ----------------------
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if hasattr(shape, "text") and "3 可行性说明" in shape.text:
        tf = shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = "3 可行性说明"
        p.font.size = Pt(28)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "\n科研研究基础："
        p.font.size = Pt(18)
        p = tf.add_paragraph()
        p.text = "本单位作为主要承担单位，参与了国铁集团重点课题 N2023S005《铁路客站管控平台与既有旅客服务系统深化融合技术研究》，2025年通过验收，成果已在多个车站试点应用。"
        p.font.size = Pt(16)
        p.level = 1
        p = tf.add_paragraph()
        p.text = "形成的相关成果："
        p.font.size = Pt(16)
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• 专利：一种车站引导屏的闭环管理方法和系统、一种基于自适应策略的引导屏车次组合排序展示方法"
        p.font.size = Pt(14)
        p.level = 2
        p = tf.add_paragraph()
        p.text = "• 软著：《铁路旅客服务与生产管控平台旅客服务软件》V2.0"
        p.font.size = Pt(14)
        p.level = 2
        p = tf.add_paragraph()
        p.text = "• 论文：多篇SCI及核心期刊论文"
        p.font.size = Pt(14)
        p.level = 2
        
        p = tf.add_paragraph()
        p.text = "\n产品应用情况："
        p.font.size = Pt(18)
        p = tf.add_paragraph()
        p.text = "系统自2025年起在济南西站、西安北站试点，2025年底在重庆东站正式投入使用，覆盖多类显示终端，运行稳定，为标准修订提供了实际运行数据和应用反馈。"
        p.font.size = Pt(16)
        p.level = 1

# ----------------------
# 第7页：标准章节结构及内容说明
# ----------------------
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if hasattr(shape, "text") and "4 标准章节结构及内容说明" in shape.text:
        tf = shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = "4 标准章节结构及内容说明"
        p.font.size = Pt(28)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "\n本次标准修订主要内容："
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = "1. 引导屏控制卡回读功能标准化：新增回读功能技术要求，实现终端设备运行状态实时监控与远程诊断"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "2. 加密审计机制引入：在系统通信、数据存储、信息展示等关键环节引入加密与审计机制，确保内容安全可追溯"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "3. 全彩显示屏版式规范统一：制定信息展示版式规范，涵盖字体、颜色、布局、动画效果等，提升信息可读性"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "4. 新型显示技术适配：结合Mini LED等新型显示技术和智能控制单元，修订设备性能、通信协议、控制方式等指标"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "5. 系统通信与接口标准化：明确与集成管理平台的通信协议、数据格式、接口规范"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "6. 设备性能与可靠性提升：对各类设备的性能指标、工作环境、防护等级、电磁兼容性等提出要求"
        p.font.size = Pt(16)
        p.level = 1

# ----------------------
# 第8页：工作进度计划
# ----------------------
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if hasattr(shape, "text") and "5 标准制定工作进度计划" in shape.text:
        tf = shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = "5 标准制定工作进度计划"
        p.font.size = Pt(28)
        p.font.bold = True
        
        current_year = datetime.datetime.now().year
        
        p = tf.add_paragraph()
        p.text = f"\n• {current_year}年1-3月：成立标准编制组，开展需求调研和资料收集"
        p.font.size = Pt(18)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"• {current_year}年4-6月：完成标准草案编写，开展内部评审"
        p.font.size = Pt(18)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"• {current_year}年7-9月：广泛征求意见，修改完善形成标准送审稿"
        p.font.size = Pt(18)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"• {current_year}年10-12月：通过专家审查，形成标准报批稿"
        p.font.size = Pt(18)
        p.level = 1

# 保存生成的PPT
prs.save(output_path)
print(f"PPT已成功生成并保存到: {output_path}")
print("所有内容已按照模板要求完成填充！")
